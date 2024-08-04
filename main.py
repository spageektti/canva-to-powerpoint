import os
import argparse
from pptx import Presentation
from pptx.util import Inches

# function to convert pixels to inches
def pixels_to_inches(pixels):
    return pixels / 96.0  # 1 inch = 96 pixels

# function to create the PowerPoint presentation
def create_presentation(gif_directory, output_path, width, height):
    # create a new PowerPoint presentation
    presentation = Presentation()

    # set the slide dimensions
    presentation.slide_width = Inches(width)
    presentation.slide_height = Inches(height)

    print(f"Creating PowerPoint presentation with dimensions {width*96}x{height*96}.")  # Convert inches to pixels (1 inch = 96 pixels)
    print(f"Reading GIF files from: {gif_directory}")

    # loop through the files in the directory
    gif_files = [f for f in os.listdir(gif_directory) if f.endswith('.gif')]
    gif_files.sort(key=lambda f: int(f.split('.')[0]))  # ensure correct order

    for gif_file in gif_files:
        gif_path = os.path.join(gif_directory, gif_file)
        if os.path.isfile(gif_path):
            print(f"Adding {gif_file} to slide.")
            slide_layout = presentation.slide_layouts[5]  # 5 corresponds to a blank slide layout
            slide = presentation.slides.add_slide(slide_layout)
            
            # add the GIF to the slide
            left = Inches(0)
            top = Inches(0)
            width = presentation.slide_width
            height = presentation.slide_height
            slide.shapes.add_picture(gif_path, left, top, width, height)

    # save the presentation
    presentation.save(output_path)
    print(f"Presentation saved as {output_path}")

# setup command-line argument parsing
parser = argparse.ArgumentParser(description='Create a PowerPoint presentation from Canva using GIF files.')
parser.add_argument('--gif_dir', type=str, default=os.getcwd(), help='Directory containing the GIF files.')
parser.add_argument('--output_path', type=str, default=os.path.join(os.getcwd(), 'output.pptx'), help='Path to save the output PowerPoint file.')
parser.add_argument('--width', type=float, help='Width of the slides in inches.')
parser.add_argument('--height', type=float, help='Height of the slides in inches.')
parser.add_argument('--width_px', type=int, help='Width of the slides in pixels.')
parser.add_argument('--height_px', type=int, help='Height of the slides in pixels.')

args = parser.parse_args()

# prioritize width_px and height_px if provided, otherwise use width and height
if args.width_px:
    width_in = pixels_to_inches(args.width_px)
else:
    width_in = args.width if args.width else 20  # default width in inches

if args.height_px:
    height_in = pixels_to_inches(args.height_px)
else:
    height_in = args.height if args.height else 11.25  # default height in inches

# call the function with parsed arguments
create_presentation(args.gif_dir, args.output_path, width_in, height_in)
